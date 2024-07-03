import yaml
from pptx import Presentation
from pptx.util import Inches

def yaml_to_pptx(yaml_file, pptx_file):
    with open(yaml_file, 'r') as file:
        data = yaml.safe_load(file)

    prs = Presentation()

    for slide_data in data['slides']:
        slide_layout = prs.slide_layouts[slide_data.get('layout', 0)]
        slide = prs.slides.add_slide(slide_layout)

        if 'title' in slide_data:
            title = slide.shapes.title
            title.text = slide_data['title']

        if 'content' in slide_data:
            body = slide.placeholders[1]
            for item in slide_data['content']:
                p = body.text_frame.add_paragraph()
                p.text = item

        if 'images' in slide_data:
            for idx, image_path in enumerate(slide_data['images']):
                left = Inches(1)
                top = Inches(2 + idx * 2)
                slide.shapes.add_picture(image_path, left, top)

    prs.save(pptx_file)

yaml_to_pptx('slides.yaml', 'presentation.pptx')
